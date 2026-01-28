# -*- coding: utf-8 -*-
"""
Скрипт для создания презентации в формате PowerPoint
на основе научной статьи и технической документации
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """Создание презентации PowerPoint"""
    
    # Создаем презентацию
    prs = Presentation()
    
    # Устанавливаем размер слайда (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== СЛАЙД 1: ТИТУЛЬНЫЙ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Титульный слайд
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Применение искусственного интеллекта (ИИ)\nпри выборе композита для реставраций\nжевательных зубов"
    subtitle.text = "Техническая часть разработки\nComposeAI — ИИ-система выбора композита"
    
    # Форматирование заголовка
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # ========== СЛАЙД 2: ЦЕЛЬ ИССЛЕДОВАНИЯ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Заголовок и содержимое
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Цель исследования"
    tf = content.text_frame
    tf.text = "Создать автоматизированную ИИ-систему, которая:"
    
    p = tf.add_paragraph()
    p.text = "• Нормализует ЭМГ-данные с разных аппаратов"
    p.level = 1
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Анализирует научные статьи и извлекает рекомендации"
    p.level = 1
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Обучается на клинических данных из литературы"
    p.level = 1
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Рекомендует оптимальные композиты на основе комплексного анализа"
    p.level = 1
    
    # ========== СЛАЙД 3: ТЕХНИЧЕСКИЙ СТЕК ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Технический стек"
    tf = content.text_frame
    tf.text = "Язык программирования и библиотеки:"
    
    p = tf.add_paragraph()
    p.text = "• Python 3.9+ — основной язык разработки"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Streamlit — веб-фреймворк для интерактивных приложений"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• scikit-learn — машинное обучение (Random Forest, Gradient Boosting)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• pandas, numpy — обработка данных"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• PyPDF2, pdfplumber — извлечение текста из PDF"
    p.level = 1
    
    # ========== СЛАЙД 4: АРХИТЕКТУРА СИСТЕМЫ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Архитектура системы"
    tf = content.text_frame
    tf.text = "Модульная структура:"
    
    p = tf.add_paragraph()
    p.text = "• app.py — основной веб-интерфейс (Streamlit)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Код_нормализации_ЭМГ.py — нормализация ЭМГ-данных"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• composite_selector.py — логика выбора композита"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• knowledge_extractor.py — извлечение знаний из статей"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• model_trainer.py — обучение ML-модели"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• База_композитов.json — база данных композитов (180+ материалов)"
    p.level = 1
    
    # ========== СЛАЙД 5: НОРМАЛИЗАЦИЯ ЭМГ (КОД) ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Пустой слайд
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Нормализация ЭМГ-данных"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    
    # Добавляем код
    code_text = """class EMGNormalizer:
    CONTROL_VALUES_SYNAPSYS = {
        MeasurementCondition.CHEWING: {
            MuscleType.MASSETER: 352.5,
            MuscleType.TEMPORALIS: 224.0,
        }
    }
    
    def normalize(self, value, apparatus, 
                  muscle, condition):
        # Приведение к единой шкале
        control_synapsys = self.CONTROL_VALUES_SYNAPSYS
        control_apparatus = self.get_control(apparatus)
        coefficient = (control_synapsys[condition][muscle] / 
                      control_apparatus[condition][muscle])
        return value * coefficient"""
    
    code_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    code_frame = code_shape.text_frame
    code_frame.text = code_text
    code_frame.word_wrap = True
    
    for paragraph in code_frame.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(0, 0, 139)
    
    # ========== СЛАЙД 6: ВЫБОР КОМПОЗИТА (КОД) ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Алгоритм выбора композита"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    
    code_text = """def calculate_composite_score(
    self, composite, emg_features, patient_data):
    score = 0.0
    weights = {
        'microhardness': 0.3,
        'wear_resistance': 0.25,
        'polymerization_shrinkage': 0.2,
        'filler_content': 0.15
    }
    
    # Правило 1: Усадка ≤ 3%
    if composite['shrinkage'] > 3:
        return 0  # Исключаем
    
    # Правило 2: Наполнитель 25-50%
    filler = composite['filler_content']
    if 25 <= filler < 50:
        score += 0.15  # Бонус
    
    # Взвешенная оценка характеристик
    score += weights['microhardness'] * hardness_score
    score += weights['wear_resistance'] * wear_score
    return score"""
    
    code_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    code_frame = code_shape.text_frame
    code_frame.text = code_text
    code_frame.word_wrap = True
    
    for paragraph in code_frame.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = RGBColor(0, 0, 139)
    
    # ========== СЛАЙД 7: МАШИННОЕ ОБУЧЕНИЕ (КОД) ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Модуль машинного обучения"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    
    code_text = """from sklearn.ensemble import (
    RandomForestClassifier, 
    GradientBoostingClassifier,
    VotingClassifier
)
from sklearn.svm import SVC
from sklearn.neighbors import KNeighborsClassifier

# Ансамбль моделей
ensemble = VotingClassifier(
    estimators=[
        ('rf', RandomForestClassifier(
            n_estimators=200,
            class_weight='balanced'
        )),
        ('gb', GradientBoostingClassifier(
            n_estimators=200,
            learning_rate=0.1
        )),
        ('svm', SVC(probability=True)),
        ('knn', KNeighborsClassifier())
    ],
    voting='soft'
)

# Обучение
ensemble.fit(X_train, y_train)
accuracy = ensemble.score(X_test, y_test)"""
    
    code_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    code_frame = code_shape.text_frame
    code_frame.text = code_text
    code_frame.word_wrap = True
    
    for paragraph in code_frame.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = RGBColor(0, 0, 139)
    
    # ========== СЛАЙД 8: ВЕБ-ИНТЕРФЕЙС (КОД) ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Веб-интерфейс (Streamlit)"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    
    code_text = """import streamlit as st

st.title("🦷 ComposeAI")

# Ввод ЭМГ-данных
masseter_right = st.number_input(
    "Жевательная мышца правая (мкВ)",
    min_value=0.0
)

# Выбор локализации
localization = st.selectbox(
    "📍 Локализация",
    ["окклюзионная поверхность",
     "апроксимальная(-ые) поверхность(-ти)"]
)

# Кнопка анализа
if st.button("Найти оптимальный композит"):
    results = selector.select_composite(patient_data)
    st.success(f"Рекомендуется: {results[0].name}")"""
    
    code_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    code_frame = code_shape.text_frame
    code_frame.text = code_text
    code_frame.word_wrap = True
    
    for paragraph in code_frame.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = RGBColor(0, 0, 139)
    
    # ========== СЛАЙД 9: БАЗА ДАННЫХ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "База данных композитов"
    tf = content.text_frame
    tf.text = "Характеристики композитов (180+ материалов):"
    
    p = tf.add_paragraph()
    p.text = "• Микротвердость (KHN — число твердости по Кнупу)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Полимеризационная усадка (%)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Содержание наполнителя (%)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Износостойкость (низкая, средняя, высокая, очень высокая)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Глубина полимеризации (мм)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Производители: 3M ESPE, Dentsply Sirona, Ivoclar Vivadent, VOCO и др."
    p.level = 1
    
    # ========== СЛАЙД 10: ПРАВИЛА ВЫБОРА ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Правила выбора композита"
    tf = content.text_frame
    tf.text = "Алгоритм основан на научных публикациях:"
    
    p = tf.add_paragraph()
    p.text = "Правило 1: Исключение композитов с усадкой > 3%"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "   Высокая усадка → нарушение краевого прилегания → вторичный кариес"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Правило 2: Оптимальный наполнитель 25-50%"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "   Лучшее сочетание прочности и устойчивости к износу"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Ранжирование:"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "   • Приоритетный вариант: усадка ≤3% И наполнитель 25-50%"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "   • Альтернативный вариант: усадка ≤3% И наполнитель 55-70%"
    p.level = 2
    
    # ========== СЛАЙД 11: ИЗВЛЕЧЕНИЕ ЗНАНИЙ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Извлечение знаний из статей"
    tf = content.text_frame
    tf.text = "Модуль KnowledgeExtractor анализирует:"
    
    p = tf.add_paragraph()
    p.text = "• Рекомендации по использованию конкретных композитов"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Интерпретацию ЭМГ-показателей"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Клинические критерии выбора материалов"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Технические характеристики композитов"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Пары 'ЭМГ-данные → рекомендованный композит'"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Формат хранения: JSON (JavaScript Object Notation)"
    p.level = 1
    p.font.italic = True
    
    # ========== СЛАЙД 12: ОБУЧЕНИЕ МОДЕЛИ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Обучение модели машинного обучения"
    tf = content.text_frame
    tf.text = "Процесс обучения:"
    
    p = tf.add_paragraph()
    p.text = "1. Автоматическое извлечение пар 'ЭМГ → композит' из статей"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Генерация синтетических данных на основе реальных паттернов"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Расширенный поиск статей (PubMed, arXiv)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Использование ансамбля моделей (RF + GB + SVM + KNN)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. Балансировка классов для равномерного представления композитов"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "6. Автоматическое обучение при запуске приложения"
    p.level = 1
    
    # ========== СЛАЙД 13: РЕЗУЛЬТАТЫ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Результаты работы системы"
    tf = content.text_frame
    tf.text = "ИИ-система предоставляет:"
    
    p = tf.add_paragraph()
    p.text = "• Карточки композитов с метриками"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Категоризацию: 'Приоритетный вариант' / 'Альтернативный вариант'"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Обоснование выбора на основе научных данных"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Дополнительные фильтры (страна, производитель, год, цена)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Топ-5 лучших вариантов с оценками"
    p.level = 1
    
    # ========== СЛАЙД 14: РАЗВЕРТЫВАНИЕ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Развертывание системы"
    tf = content.text_frame
    tf.text = "Облачная платформа:"
    
    p = tf.add_paragraph()
    p.text = "• Streamlit Cloud — облачный сервис для веб-приложений"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• GitHub — система контроля версий и автоматическое обновление"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Кэширование данных для ускорения работы"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Ленивая загрузка модели (только при необходимости)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Доступность на всех платформах через веб-браузер"
    p.level = 1
    
    # ========== СЛАЙД 15: ВЫВОДЫ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Выводы"
    tf = content.text_frame
    tf.text = "ИИ-система ComposeAI:"
    
    p = tf.add_paragraph()
    p.text = "✓ Автоматизирует выбор композита на основе объективных данных"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Использует актуальные научные данные из литературы"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Обучается на клинических данных и улучшается со временем"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Предоставляет обоснованные рекомендации с объяснением"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "⚠ ИИ — инструмент поддержки клинического решения,"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  а не замена клинического опыта врача"
    p.level = 2
    
    # Сохраняем презентацию
    output_path = os.path.join(os.path.dirname(__file__), "Презентация_ComposeAI.pptx")
    prs.save(output_path)
    print(f"✅ Презентация создана: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Ошибка: Необходимо установить библиотеку python-pptx")
        print("Установите: pip install python-pptx")
    except Exception as e:
        print(f"❌ Ошибка при создании презентации: {e}")
        import traceback
        traceback.print_exc()
